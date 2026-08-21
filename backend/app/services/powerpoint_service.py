from pathlib import Path
from io import BytesIO
import requests
from PIL import Image, ImageOps, ImageDraw, ImageFilter



try:
    import pytesseract
except Exception:
    pytesseract = None

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


GENERATED_DIR = Path("generated_files")
GENERATED_DIR.mkdir(exist_ok=True)

SW = Inches(13.333)
SH = Inches(7.5)


# ============================================================
# DESIGN SYSTEM
# ============================================================

def rgb(hex_value):
    h = hex_value.replace("#", "")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def theme_colors(theme=None):
    t = theme or {}
    return {
        "navy": rgb(t.get("background", "#071321")),
        "navy2": rgb(t.get("surface", "#10243A")),
        "blue": rgb(t.get("accent", "#0877E8")),
        "orange": rgb(t.get("secondary_accent", "#FF9F1C")),
        "white": rgb("#FFFFFF"),
        "paper": rgb("#F5F8FC"),
        "ink": rgb("#0E1A2B"),
        "muted": rgb("#64748B"),
        "line": rgb("#D8E0EA"),
    }


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def shape(slide, kind, x, y, w, h, fill, radius=False, transparency=0):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else kind
    s = slide.shapes.add_shape(st, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.fill.transparency = transparency
    s.line.fill.background()
    return s


def circle(slide, x, y, d, fill, transparency=0):
    return shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill, False, transparency)


def text(slide, value, x, y, w, h, size=18, color=None,
         bold=False, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(value)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return box


def add_rule(slide, x, y, w, color):
    shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.045), color)


def page(slide, n, color):
    text(slide, f"{n:02d}", Inches(12.15), Inches(6.85),
         Inches(0.45), Inches(0.25), 10, color, False,
         PP_ALIGN.RIGHT)


# ============================================================
# IMAGE HANDLING
# ============================================================

BAD_IMAGE_WORDS = (
    "logo",
    "screenshot",
    "infographic",
    "diagram",
    "template",
    "presentation",
    "powerpoint",
    "slide",
    "chart",
    "graph",
    "poster",
    "meme",
    "watermark",
    "collage",
    "thumbnail",
    "banner",
    "flyer",
    "advertisement",
    "ad",
)


def image_url(item):
    if isinstance(item, str):
        return item
    if isinstance(item, dict):
        return item.get("url") or item.get("image_url")
    return None


def image_description(item):
    if isinstance(item, dict):
        return str(
            item.get("description")
            or item.get("title")
            or item.get("alt")
            or ""
        ).lower()
    return ""


def good_image(item):
    """Reject screenshots/infographics/text-heavy images before PPT insertion."""
    url = image_url(item)
    if not url:
        return False

    description = image_description(item)
    if any(word in description for word in BAD_IMAGE_WORDS):
        return False

    return True


def image_is_visual(url):
    """Validate the actual pixels, not only the web-result description."""
    try:
        r = requests.get(
            url,
            timeout=15,
            headers={"User-Agent": "Mozilla/5.0"},
        )
        r.raise_for_status()
        im = Image.open(BytesIO(r.content)).convert("RGB")

        # Avoid tiny thumbnails / icons.
        if im.width < 700 or im.height < 450:
            return False

        # Reject images that are mostly text, such as screenshots and infographics.
        if pytesseract is not None:
            sample = ImageOps.contain(im, (1100, 700))
            text = pytesseract.image_to_string(sample, config="--psm 6")
            words = len(text.split())
            if words >= 45:
                return False

        return True
    except Exception:
        # If validation cannot inspect the image, let the normal downloader decide.
        return True


def download_png(url):
    try:
        r = requests.get(
            url,
            timeout=20,
            headers={
                "User-Agent": "Mozilla/5.0",
                "Accept": "image/avif,image/webp,image/apng,image/svg+xml,image/*,*/*;q=0.8",
            },
        )
        r.raise_for_status()

        source = BytesIO(r.content)
        im = Image.open(source)
        im = ImageOps.exif_transpose(im)

        # Always normalize web images to PNG.
        if im.mode not in ("RGB", "RGBA"):
            im = im.convert("RGBA")

        out = BytesIO()
        im.save(out, "PNG")
        out.seek(0)
        return out
    except Exception as e:
        print(f"Image download failed: {e}")
        return None


def prepare_crop(url, width=1200, height=800):
    stream = download_png(url)
    if not stream:
        return None

    try:
        im = Image.open(stream).convert("RGB")
        im = ImageOps.fit(
            im,
            (width, height),
            method=Image.Resampling.LANCZOS,
            centering=(0.5, 0.5),
        )

        # Slight sharpening after resizing.
        im = im.filter(ImageFilter.UnsharpMask(radius=1, percent=110, threshold=3))

        out = BytesIO()
        im.save(out, "PNG", optimize=True)
        out.seek(0)
        return out
    except Exception as e:
        print(f"Image preparation failed: {e}")
        return None


def add_image(slide, url, x, y, w, h, rounded=True):
    stream = prepare_crop(url)
    if not stream:
        return False

    try:
        # Rounded image is represented by a clean image card.
        if rounded:
            shape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, w, h, rgb("#FFFFFF"), True
            )
        slide.shapes.add_picture(stream, x, y, width=w, height=h)
        return True
    except Exception as e:
        print(f"Could not insert image: {e}")
        return False


def add_circle_image(slide, url, x, y, d):
    stream = prepare_crop(url, 1000, 1000)
    if not stream:
        return False

    try:
        im = Image.open(stream).convert("RGBA")
        im = ImageOps.fit(
            im, (900, 900),
            method=Image.Resampling.LANCZOS,
            centering=(0.5, 0.5)
        )

        mask = Image.new("L", im.size, 0)
        draw = ImageDraw.Draw(mask)
        draw.ellipse((0, 0, im.width - 1, im.height - 1), fill=255)
        im.putalpha(mask)

        out = BytesIO()
        im.save(out, "PNG")
        out.seek(0)

        slide.shapes.add_picture(out, x, y, width=d, height=d)
        return True
    except Exception as e:
        print(f"Could not insert circular image: {e}")
        return False


def choose_image(slide_data, images, used_urls=None):
    """
    Select a unique high-quality image.

    IMPORTANT:
    - Never reuse an image.
    - Reject screenshots, infographics and text-heavy images.
    - If no unused suitable image exists, return None.
    """
    used_urls = used_urls if used_urls is not None else set()

    candidates = []

    idx = slide_data.get("image_index")

    if isinstance(idx, int) and 0 <= idx < len(images):
        candidates.append(images[idx])

    candidates.extend(images)

    seen = set()

    for item in candidates:
        url = image_url(item)

        if not url:
            continue

        # Don't inspect the same URL repeatedly.
        if url in seen:
            continue

        seen.add(url)

        # NEVER reuse an image.
        if url in used_urls:
            continue

        # Reject bad search results.
        if not good_image(item):
            continue

        # Validate actual image.
        if not image_is_visual(url):
            continue

        # Reserve this image permanently.
        used_urls.add(url)

        print(f">>> Selected unique image: {url}")

        return url

    print(">>> No unused suitable image available for this slide.")

    # IMPORTANT:
    # Do NOT reuse an old image.
    return None


# ============================================================
# SLIDE LAYOUTS — INSPIRED BY THE REFERENCE DESIGN
# ============================================================

def hero(prs, title, theme, images, used_urls=None):
    c = theme_colors(theme)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, c["navy"])

    # Editorial blue/orange geometry
    circle(s, Inches(9.0), Inches(4.9), Inches(3.2), c["blue"], 8)
    circle(s, Inches(10.8), Inches(0.3), Inches(0.5), c["blue"])
    circle(s, Inches(11.5), Inches(5.7), Inches(0.8), c["orange"])
    circle(s, Inches(9.55), Inches(6.15), Inches(0.38), c["blue"])
    circle(s, Inches(8.8), Inches(6.0), Inches(0.55), c["orange"])

    text(s, "MAK-AI", Inches(0.7), Inches(0.55),
         Inches(2), Inches(0.3), 14, c["white"], True)

    # Use a large, clean photographic image.
    url = choose_image({"image_index": 0}, images, used_urls)
    if url:
        add_circle_image(s, url, Inches(8.8), Inches(0.55), Inches(3.4))

    text(s, title, Inches(0.7), Inches(2.05),
         Inches(7.4), Inches(1.8), 38, c["white"], True)

    # Blue accent word/line
    add_rule(s, Inches(0.72), Inches(4.35), Inches(0.75), c["orange"])

    text(s, "A visual journey through the topic",
         Inches(0.72), Inches(4.55), Inches(5.5), Inches(0.45),
         16, c["white"])

    text(s, "AI-generated presentation",
         Inches(0.72), Inches(6.65), Inches(3.5), Inches(0.3),
         10, c["muted"])


def overview(s, data, c, n):
    bg(s, c["paper"])

    text(s, data.get("title", "Overview"), Inches(0.75), Inches(0.55),
         Inches(8), Inches(0.55), 29, c["ink"], True)
    add_rule(s, Inches(0.75), Inches(1.35), Inches(0.45), c["blue"])
    circle(s, Inches(11.85), Inches(0.5), Inches(0.5), c["blue"])
    text(s, f"{n:02d}", Inches(11.95), Inches(0.66),
         Inches(0.3), Inches(0.2), 10, c["white"], True, PP_ALIGN.CENTER)

    items = data.get("content", [])[:4]
    xs = [0.85, 3.85, 6.85, 9.85]
    colors = [c["blue"], c["orange"], c["blue"], c["orange"]]

    for i, item in enumerate(items):
        text(s, f"{i+1:02d}", Inches(xs[i]), Inches(2.0),
             Inches(0.5), Inches(0.3), 18, colors[i], True)
        add_rule(s, Inches(xs[i]), Inches(2.42), Inches(0.4), colors[i])
        circle(s, Inches(xs[i]), Inches(2.8), Inches(0.62), colors[i])
        text(s, "●", Inches(xs[i]+0.15), Inches(2.96),
             Inches(0.3), Inches(0.2), 9, c["white"], True, PP_ALIGN.CENTER)
        text(s, item, Inches(xs[i]), Inches(3.7),
             Inches(2.15), Inches(1.35), 14, c["ink"], True)

    # subtle decorative dots
    for i in range(4):
        circle(s, Inches(0.8+i*0.18), Inches(6.65), Inches(0.06), c["blue"])


def image_text(s, data, c, n, images, used_urls=None):
    bg(s, c["paper"])

    text(s, data.get("title", "Key Concept"), Inches(0.75), Inches(0.55),
         Inches(8), Inches(0.7), 29, c["ink"], True)
    add_rule(s, Inches(0.75), Inches(1.35), Inches(0.55), c["blue"])
    page(s, n, c["muted"])

    url = choose_image(data, images, used_urls)

    # Large image block
    if url:
        add_image(s, url, Inches(7.65), Inches(1.15),
                  Inches(4.75), Inches(4.95), True)
    else:
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.65), Inches(1.15),
              Inches(4.75), Inches(4.95), c["navy2"], True)

    items = data.get("content", [])[:4]
    y = 1.85

    for item in items:
        circle(s, Inches(0.8), Inches(y+0.04), Inches(0.16), c["blue"])
        text(s, item, Inches(1.1), Inches(y),
             Inches(5.75), Inches(0.65), 16, c["ink"])
        y += 0.82

    circle(s, Inches(11.9), Inches(5.95), Inches(0.55), c["orange"])


def timeline(s, data, c, n, images, used_urls=None):
    bg(s, c["paper"])

    text(s, data.get("title", "Timeline"), Inches(0.75), Inches(0.55),
         Inches(8), Inches(0.65), 29, c["ink"], True)
    page(s, n, c["muted"])

    items = data.get("content", [])[:4]
    x_positions = [1.05, 4.05, 7.05, 10.05]

    shape(s, MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.15),
          Inches(10.9), Inches(0.055), c["line"])

    for i, item in enumerate(items):
        x = x_positions[i]
        color = c["blue"] if i % 2 == 0 else c["orange"]

        url = choose_image(
            {"image_index": data.get("image_indices", [None]*4)[i]
             if isinstance(data.get("image_indices"), list)
             and i < len(data.get("image_indices")) else None},
            images,
            used_urls,
        )

        if url:
            add_circle_image(s, url, Inches(x-0.35), Inches(1.55), Inches(1.7))
        else:
            circle(s, Inches(x-0.35), Inches(1.55), Inches(1.7), c["navy2"])

        circle(s, Inches(x+0.35), Inches(3.93), Inches(0.42), color)

        text(s, f"{i+1:02d}", Inches(x-0.15), Inches(4.55),
             Inches(0.3), Inches(0.25), 12, color, True, PP_ALIGN.CENTER)

        text(s, item, Inches(x-0.7), Inches(5.0),
             Inches(1.95), Inches(1.0), 13, c["ink"], True, PP_ALIGN.CENTER)


def dark_cards(s, data, c, n, images, used_urls=None):
    bg(s, c["navy"])

    circle(s, Inches(11.4), Inches(-0.55), Inches(2.0), c["blue"], 5)
    circle(s, Inches(0.25), Inches(6.75), Inches(0.5), c["orange"])

    text(s, f"{n:02d}", Inches(0.75), Inches(0.55),
         Inches(0.5), Inches(0.3), 12, c["blue"], True)
    text(s, data.get("title", "Key Topics"), Inches(0.75), Inches(1.0),
         Inches(8.5), Inches(0.7), 30, c["white"], True)
    add_rule(s, Inches(0.75), Inches(1.85), Inches(0.55), c["blue"])

    items = data.get("content", [])[:5]
    x = 0.7
    card_w = 2.35

    for i, item in enumerate(items):
        x = 0.65 + i * 2.48
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(x), Inches(2.45), Inches(card_w), Inches(3.65),
              c["navy2"], True)

        # Use a different high-quality research image for each card.
        card_url = None

        if i < 3:
            card_url = choose_image(
                {"image_index": i},
                images,
                used_urls
            )
        if card_url:
            add_image(
                s, card_url,
                Inches(x+0.12), Inches(2.6),
                Inches(2.11), Inches(1.15), True
            )

        text(s, f"{i+1:02d}", Inches(x+0.18), Inches(4.0),
             Inches(0.45), Inches(0.3), 13, c["orange"], True)
        text(s, item, Inches(x+0.18), Inches(4.45),
             Inches(1.95), Inches(1.2), 13, c["white"], True)


def stats(s, data, c, n):
    bg(s, c["paper"])

    text(s, data.get("title", "Highlights"), Inches(0.75), Inches(0.55),
         Inches(9), Inches(0.65), 29, c["ink"], True)
    page(s, n, c["muted"])

    items = data.get("content", [])[:3]
    colors = [c["blue"], c["orange"], c["blue"]]

    for i, item in enumerate(items):
        x = 0.9 + i * 4.1
        circle(s, Inches(x+1.15), Inches(1.75), Inches(1.0), colors[i])
        text(s, f"{i+1:02d}", Inches(x+1.43), Inches(2.1),
             Inches(0.45), Inches(0.25), 12, c["white"], True, PP_ALIGN.CENTER)

        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(x), Inches(3.05), Inches(3.4), Inches(2.2),
              c["white"], True)

        text(s, item, Inches(x+0.28), Inches(3.45),
             Inches(2.85), Inches(1.25), 16, c["ink"], True, PP_ALIGN.CENTER)


def conclusion(s, data, c, n, images, used_urls=None):
    bg(s, c["navy"])

    text(s, f"{n:02d}", Inches(0.75), Inches(0.55),
         Inches(0.5), Inches(0.3), 12, c["blue"], True)

    text(s, data.get("title", "Conclusion"), Inches(0.75), Inches(1.05),
         Inches(7.2), Inches(0.7), 32, c["white"], True)
    add_rule(s, Inches(0.75), Inches(1.9), Inches(0.55), c["blue"])

    items = data.get("content", [])[:4]
    y = 2.35
    for item in items:
        circle(s, Inches(0.78), Inches(y+0.06), Inches(0.13), c["blue"])
        text(s, item, Inches(1.05), Inches(y),
             Inches(5.7), Inches(0.65), 16, c["white"])
        y += 0.8

    url = choose_image(data, images, used_urls)
    if url:
        add_circle_image(s, url, Inches(8.35), Inches(1.25), Inches(3.9))

    circle(s, Inches(10.2), Inches(5.75), Inches(0.85), c["blue"])
    circle(s, Inches(11.35), Inches(6.0), Inches(0.48), c["orange"])

    text(s, "MAK-AI  |  Thank You!", Inches(0.75), Inches(6.7),
         Inches(3.5), Inches(0.3), 11, c["muted"])


# ============================================================
# MAIN RENDERER
# ============================================================

def create_powerpoint(title, slides, theme=None, images=None):
    images = images or []
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # Intentional visual rhythm matching the Mak-AI reference:
    # dark / light / dark / light ... rather than random backgrounds.
    used_urls = set()
    hero(prs, title, theme, images, used_urls)
    c = theme_colors(theme)

    for i, data in enumerate(slides):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        layout = data.get("layout", "image_text")
        n = i + 2

        if layout == "overview":
            overview(s, data, c, n)
        elif layout in ("image_text", "split"):
            image_text(s, data, c, n, images, used_urls)
        elif layout == "timeline":
            timeline(s, data, c, n, images, used_urls)
        elif layout in ("dark_section", "cards"):
            dark_cards(s, data, c, n, images, used_urls)
        elif layout == "stats":
            stats(s, data, c, n)
        elif layout == "conclusion":
            conclusion(s, data, c, n, images, used_urls)
        else:
            image_text(s, data, c, n, images, used_urls)

    safe = "".join(
        ch if ch.isalnum() or ch in " _-" else "_"
        for ch in title
    ).strip()[:80]

    path = GENERATED_DIR / f"{safe}.pptx"
    prs.save(path)
    print(f"PowerPoint generated successfully: {path}")
    return str(path)